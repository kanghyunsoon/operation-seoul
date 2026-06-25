from html import escape
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
SOURCE = DOCS / "09_PRESENTATION_DRAFT.md"
PPTX_OUT = DOCS / "OperationSeoul_Final_Presentation.pptx"
HTML_OUT = DOCS / "OperationSeoul_Final_Presentation.html"


def parse_slides(markdown: str):
    parts = re.split(r"\n## Slide\s+\d+\.\s+", markdown)
    slides = []
    for part in parts[1:]:
        lines = [line.rstrip() for line in part.strip().splitlines()]
        if not lines:
            continue
        title = lines[0].strip()
        body = "\n".join(lines[1:]).strip()
        slides.append((title, body))
    return slides


def clean_bullet(line: str):
    return re.sub(r"^[-*]\s+", "", line).strip()


def add_textbox(slide, left, top, width, height, text, size=22, bold=False, color=(28, 38, 55)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*color)
    return box


def build_pptx(slides):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, (title, body) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(248, 251, 255)

        add_textbox(slide, Inches(0.65), Inches(0.45), Inches(11.6), Inches(0.7), title, size=31, bold=True)
        add_textbox(slide, Inches(11.7), Inches(0.55), Inches(0.8), Inches(0.35), f"{index:02d}", size=12, bold=True, color=(37, 99, 235))

        content = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(11.7), Inches(5.5))
        frame = content.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.05)
        frame.margin_right = Inches(0.05)

        frame.clear()
        body_lines = [line.strip() for line in body.splitlines() if line.strip()]
        first = True
        for line in body_lines:
            p = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            if line.startswith("- "):
                p.text = clean_bullet(line)
                p.level = 0
                p.font.size = Pt(22)
            elif re.match(r"^\d+\.\s+", line):
                p.text = line
                p.font.size = Pt(21)
            else:
                p.text = line
                p.font.size = Pt(23 if index == 1 else 21)
            p.font.name = "Malgun Gothic"
            p.font.color.rgb = RGBColor(34, 47, 62)
            p.space_after = Pt(9)

        footer = slide.shapes.add_textbox(Inches(0.65), Inches(7.0), Inches(12), Inches(0.25))
        fp = footer.text_frame.paragraphs[0]
        fp.text = "Operation KOREA Final Submission"
        fp.alignment = PP_ALIGN.RIGHT
        fp.font.size = Pt(9)
        fp.font.color.rgb = RGBColor(100, 116, 139)

    prs.save(PPTX_OUT)


def body_to_html(body: str):
    lines = [line.rstrip() for line in body.splitlines() if line.strip()]
    html = []
    in_list = False
    for line in lines:
        if line.startswith("- "):
            if not in_list:
                html.append("<ul>")
                in_list = True
            html.append(f"<li>{escape(clean_bullet(line))}</li>")
        else:
            if in_list:
                html.append("</ul>")
                in_list = False
            html.append(f"<p>{escape(line)}</p>")
    if in_list:
        html.append("</ul>")
    return "\n".join(html)


def build_html(slides):
    sections = []
    for index, (title, body) in enumerate(slides, start=1):
        sections.append(
            f"""
            <section class="slide">
              <div class="num">{index:02d}</div>
              <h1>{escape(title)}</h1>
              <div class="body">{body_to_html(body)}</div>
              <footer>Operation KOREA Final Submission</footer>
            </section>
            """
        )

    HTML_OUT.write_text(
        f"""<!doctype html>
<html lang="ko">
<head>
  <meta charset="utf-8">
  <title>Operation KOREA Final Presentation</title>
  <style>
    @page {{ size: 16:9 landscape; margin: 0; }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; background: #e5edf7; color: #172033; font-family: 'Malgun Gothic', 'Noto Sans KR', Arial, sans-serif; }}
    .slide {{ position: relative; width: 100vw; height: 100vh; page-break-after: always; padding: 52px 70px; background: #f8fbff; overflow: hidden; }}
    .slide::before {{ content: ''; position: absolute; left: 0; top: 0; width: 12px; height: 100%; background: #2563eb; }}
    .num {{ position: absolute; right: 70px; top: 52px; color: #2563eb; font-weight: 800; }}
    h1 {{ margin: 0 0 44px; max-width: 980px; font-size: 42px; line-height: 1.18; }}
    .body {{ max-width: 1080px; font-size: 25px; line-height: 1.55; }}
    p {{ margin: 0 0 18px; }}
    ul {{ margin: 0; padding-left: 30px; }}
    li {{ margin: 0 0 14px; }}
    footer {{ position: absolute; right: 70px; bottom: 36px; color: #64748b; font-size: 13px; }}
  </style>
</head>
<body>
{''.join(sections)}
</body>
</html>
""",
        encoding="utf-8",
    )


def main():
    slides = parse_slides(SOURCE.read_text(encoding="utf-8"))
    build_pptx(slides)
    build_html(slides)
    print(PPTX_OUT)
    print(HTML_OUT)


if __name__ == "__main__":
    main()
